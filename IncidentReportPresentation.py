from datetime import datetime
import string
from pptx import Presentation
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.slide import Slides
from pptx.text.text import _Paragraph

from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

from UniqueColorGenerator import UniqueColorGenerator


sky_blue = RGBColor(0x00, 0x70, 0xc0)
bright_red = RGBColor(0xB2, 0x22, 0x22)


class IncidentReportPresentation:
    def __init__(self, pres: Presentation, enterprise_logo_path: string, event_table_headers, colorize_event_tables = True):
        self.prs = pres
        self.enterprise_logo = enterprise_logo_path
        self.event_table_headers = event_table_headers
        self.color_generator = UniqueColorGenerator(70)
        self.colorize_direct_causes = colorize_event_tables

    def save(self, file_path:string):
        return self.prs.save(file_path)

    def paginate_presentation(self):
        for i, slide in enumerate(self.prs.slides, start=1):
            slide_number_box = slide.shapes.add_textbox(Inches(9.45), Inches(6.90), Inches(0.5), Inches(0.5))
            slide_number_frame = slide_number_box.text_frame
            slide_number = slide_number_frame.add_paragraph()
            slide_number.text = "Page " + str(i)
            slide_number.font.size = Pt(12)
            slide_number.font.name = 'Calibri'
            slide_number.alignment = PP_ALIGN.RIGHT

        # set footer logo to slides
    def set_footer_logo_to_slides(self):
        for i, slide in enumerate(self.prs.slides, start=1):
            slide.shapes.add_picture(self.enterprise_logo, Inches(0), Inches(6.90), Inches(0.60), Inches(0.60))

    def add_textbox_to_slide(self, slide_id: int, left: float, top: float, width: float, height: float, word_wrap = True):
        slide = self.prs.slides.get(slide_id)
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        textbox.text_frame.word_wrap = word_wrap
        return textbox

    def edit_textbox(self, text_box, text = None, font_size = None, font_name:PP_ALIGN = None, color:RGBColor = None, alignment = None, bold=False):
        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]

        self.__edit_paragraph(paragraph, text, font_size, font_name, color, alignment, bold)

    def __edit_paragraph(self, paragraph:_Paragraph, text_content = None, font_size = None, font_name = None, color:RGBColor = None, alignment = None, bold=False):
        if text_content != None :
            paragraph.text = text_content
        
        if font_size != None:
            paragraph.font.size = Pt(font_size)

        if font_name != None:
            paragraph.font.name = font_name

        if color != None:    
            paragraph.font.color.rgb = color

        if alignment != None:
            paragraph.alignment = alignment

        if bold != None:
            paragraph.font.bold = bold

    def add_paragraph(self, text_box, text = '', font_size = 16 , font_name = 'Calibri', color = RGBColor(0x00, 0x00, 0x00), alignment = PP_ALIGN.LEFT, bold=False):
        text_frame = text_box.text_frame
        new_paragraph = text_frame.add_paragraph() 
        self.__edit_paragraph(new_paragraph, text, font_size, font_name, color, alignment, bold)

    def set_front_page(self, incident_site, incident_report_edition_date, incident_title, incident_title_sub_text, image_uri = None):

        if image_uri != None :
            # prs.slide_layouts[8] => picture_with_caption_layout
            selected_layout = self.prs.slide_layouts[8]
            placeholder_id = 2
        else :
            selected_layout = self.prs.slide_layouts[0]
            placeholder_id = 1

        slide1 = self.prs.slides.add_slide(selected_layout)
        title_box = self.add_textbox_to_slide(slide1.slide_id, left=2.5, top=0.10, width=5, height=0.5)

        sky_blue = RGBColor(0x00, 0x70, 0xc0)
        bright_red = RGBColor(0xB2, 0x22, 0x22)

        incident_info = incident_site + " LE " + incident_report_edition_date

        self.edit_textbox(title_box, incident_info.upper(), 18, 'Calibri', sky_blue, PP_ALIGN.CENTER, True)

        # picture_with_caption_layout: placeholders[0] => Title
        title_shape = slide1.shapes.placeholders[0]
        title = title_shape.text_frame.add_paragraph()
        self.__edit_paragraph(paragraph=title, text_content=incident_title.upper(), color=bright_red ,alignment=PP_ALIGN.CENTER, bold=True)

        # picture_with_caption_layout: placeholders[2] => Text Placeholder
        text_shape = slide1.shapes.placeholders[placeholder_id]
        sub_title = text_shape.text_frame.add_paragraph()
        self.__edit_paragraph(paragraph=sub_title, text_content=incident_title_sub_text, font_size=18, font_name='Calibri', color=bright_red ,alignment=PP_ALIGN.CENTER, bold=True)

    def set_summary_slide(self, summary_title, parts):
        # prs.slide_layouts[1] => bullet_slide_layout
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide2 = self.prs.slides.add_slide(bullet_slide_layout)

        title_slide2_shape = slide2.shapes.placeholders[0]
        title_slide2_shape.text = summary_title
        title_slide2_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xB2, 0x22, 0x22)
        title_slide2_shape.text_frame.paragraphs[0].font.bold = True

        slide2_shape_content = slide2.shapes.placeholders[1]
        slide2_shape_content_text_frame = slide2_shape_content.text_frame

        # Redimensionner la TextFrame
        slide2_shape_content.width = 5000000
        slide2_shape_content.height = 5000000

        # Déplacer la TextFrame
        slide2_shape_content.left = 2500000
        slide2_shape_content.top  = 1500000

        for i, member in enumerate(parts):
            p2 = slide2_shape_content_text_frame.add_paragraph()
            self.__edit_paragraph(p2, text_content=member, font_size=24, color=bright_red, bold=True)
            # p2.text = member
            # p2.font.color.rgb = RGBColor(0xB2, 0x22, 0x22)



    def set_context_slide(self, slide_title, context):

        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        bright_red = RGBColor(0xB2, 0x22, 0x22)

        txBox = self.add_textbox_to_slide(slide.slide_id, left=0.3, top=0, width=2, height=0.4, word_wrap=False)

        p = txBox.text_frame.paragraphs[0]
        self.__edit_paragraph(paragraph=p, text_content=slide_title, font_size=36, font_name='Calibri', color=bright_red ,alignment=PP_ALIGN.CENTER, bold=True)
        
        context = context.replace("\r", "")
        
        txBox = self.add_textbox_to_slide(slide.slide_id, 0.25, 0.5, 9.5, 6.4)
        p = txBox.text_frame.add_paragraph()
        self.__edit_paragraph(paragraph=p, text_content=context, font_size=13, font_name='Calibri', alignment=PP_ALIGN.LEFT)



    def set_team_slide(self, slide_title, team_members):
        # prs.slide_layouts[6] => blank_slide_layout
        blank_slide_layout = self.prs.slide_layouts[6]
        slide4 = self.prs.slides.add_slide(blank_slide_layout)

        txBox2 = self.add_textbox_to_slide(slide4.slide_id, 0.3, 0, 4, 0.4, False)
        tf = txBox2.text_frame
        p = tf.paragraphs[0]

        self.__edit_paragraph(p, slide_title, font_size=36, color=RGBColor(0xB2, 0x22, 0x22), bold=True)

        # Ajouter du contenu à la diapositive sous forme de deux colonnes
        txBox1 = slide4.shapes.add_textbox(Inches(0), Inches(0.5), Inches(4.5), Inches(6))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True
        tf1.auto_size = True

        txBox2 = slide4.shapes.add_textbox(Inches(4), Inches(0.5), Inches(6), Inches(6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.auto_size = True

        # Ajouter un nouveau paragraphe pour chaque membre de l'équipe
        for i, member in enumerate(team_members):
            if i < len(team_members) / 2:
                tf = tf1
            else:
                tf = tf2
            p = tf.add_paragraph()
            p.text = "• " + member  # Ajouter une puce
            p.font.size = Pt(16)
            p.level = 1  # Ajouter une puce
            p.space_before = Pt(21)


    def set_methodology_illustration(self, slide_title, illustration1 = 'Picture2.png', illustration2 = 'Picture1.png'):
        slide_layout = self.prs.slide_layouts[6]  # 6 is a blank slide
        slide7 = self.prs.slides.add_slide(slide_layout)

        title_box = self.add_textbox_to_slide(slide7.slide_id, left=2, top=0.25, width=6, height=1, word_wrap=False)
        # Add a title to the slide
        
        self.edit_textbox(title_box, text=slide_title, font_size=48, font_name='Tahoma', color=bright_red)


        # Add the images to the slide
        slide7.shapes.add_picture(illustration1, Inches(0), Inches(1), Inches(4.5), Inches(5.5))
        slide7.shapes.add_picture(illustration2, Inches(4), Inches(1.75), Inches(6), Inches(2))


        # Add text below each image
        text1_box = self.add_textbox_to_slide(slide7.slide_id, left=0.5, top=6, width=5, height=1, word_wrap=False)
        self.add_paragraph(text1_box, text="SCAT – isrs7", font_size= 44, font_name= 'Tahoma', color=bright_red)


        text2_box = self.add_textbox_to_slide(slide7.slide_id, left=4.25, top=3.5, width=6, height=1, word_wrap=False)
        self.add_paragraph(text2_box, text="ARBRE DES CAUSES", font_size=44, font_name='Tahoma', color=bright_red)


    def add_event_slide(self, event_data):
        # Ajouter une diapositive vide (la septième diapositive)
        slide_layout = self.prs.slide_layouts[6]  # 6 est une diapositive vide
        slide8 = self.prs.slides.add_slide(slide_layout)

        # Ajouter un titre à la diapositive
        title_box = self.add_textbox_to_slide(slide8.slide_id, left=0, top=0, width=10, height=1, word_wrap=True)
        title_box_content = "Evènement : " + event_data['title'] + " (" + event_data['EventType']['code'] + ")"
        self.edit_textbox(title_box, text=title_box_content, font_size=24, font_name='Calibri', color=bright_red, alignment=PP_ALIGN.CENTER, bold=True)

        # Ajouter un tableau en bas de la diapositive
        table = slide8.shapes.add_table(2, 5, Inches(0.25), Inches(1.15), Inches(9.5), Inches(10.75)).table  # Added 0.5 inch margin to the left, right and bottom


        table.columns[0].width = Inches(2.6)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(2.2)
        table.columns[3].width = Inches(1.5)
        table.columns[4].width = Inches(1)
        
        # Définir les titres des colonnes
        for i in range(len(self.event_table_headers)):
            cell = table.cell(0, i)
            cell.text_frame.auto_size = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            cell_text = cell.text_frame.paragraphs[0]  # Center column titles
            self.__edit_paragraph(cell_text, text_content=self.event_table_headers[i], font_size=13, font_name='Calibri', alignment=PP_ALIGN.CENTER, bold=True)

        table.rows[0].height = Inches(0.5)  # Définir la hauteur à 1 pouce


        for direct_cause in event_data['directCauses']:
            cell = table.cell(1, 0)
            
            direct_cause_type_ref = direct_cause['DirectCauseTypeId']['name'] + ' :'

            random_direct_cause_event_color = sky_blue

            if self.colorize_direct_causes:
                random_color = self.color_generator.generate_unique_color()
                random_direct_cause_event_color =  RGBColor(random_color[0], random_color[1], random_color[2])

            cell_text = cell.text_frame.add_paragraph()  # Center column titles
            self.__edit_paragraph(cell_text, text_content=direct_cause_type_ref, font_size=11, font_name='Calibri', color=random_direct_cause_event_color, bold=True)

            direct_cause_description = direct_cause['description']
            cell_text = cell.text_frame.add_paragraph()  # Center column titles

            direct_cause_description = direct_cause_description.replace("\r","").replace("\n","")
            self.__edit_paragraph(cell_text, text_content=direct_cause_description, font_size=11, font_name='Calibri')

            cell_text = cell.text_frame.add_paragraph()

            for root_cause in direct_cause['fundamentalCauses']:
                cell1 = table.cell(1, 1)

                root_cause_type_ref = "[" + direct_cause['DirectCauseTypeId']['code'] + "] - " + root_cause['FundamentalCauseTypeId']['name'] + " :" 
                cell_text1 = cell1.text_frame.add_paragraph()  
                self.__edit_paragraph(cell_text1, text_content=root_cause_type_ref, font_size=11, font_name='Calibri', color=random_direct_cause_event_color, bold=True)

                root_cause_description = root_cause['description']
                cell_text1 = cell1.text_frame.add_paragraph()

                root_cause_description = root_cause_description.replace("\r","").replace("\n","")
                self.__edit_paragraph(cell_text1, text_content=root_cause_description, font_size=11, font_name='Calibri')

                cell_text1 = cell1.text_frame.add_paragraph()

                for improvement_action in root_cause['improvementActions']:
                    cell2 = table.cell(1, 2)
                    
                    root_cause_type_code = improvement_action['improvementActionTypeId']['code'] 

                    improvement_action_type_ref = "[" + root_cause['FundamentalCauseTypeId']['code'] + "] - " + improvement_action['improvementActionTypeId']['process'] + " :" 
                    cell_text2 = cell2.text_frame.add_paragraph()  
                    self.__edit_paragraph(cell_text2, text_content=improvement_action_type_ref, font_size=11, font_name='Calibri', color=random_direct_cause_event_color, bold=True)

                    root_cause_description = improvement_action['suggestion']
                    cell_text2 = cell2.text_frame.add_paragraph()  
                    
                    root_cause_description = root_cause_description.replace("\r", "").replace("\n","")
                    self.__edit_paragraph(cell_text2, text_content=root_cause_description, font_size=11, font_name='Calibri')

                    cell_text2 = cell2.text_frame.add_paragraph()
                    cell_text2 = cell2.text_frame.add_paragraph()

                    cell3 = table.cell(1, 3)
                    responsibilities_ref = "[" + improvement_action['improvementActionTypeId']['code'] + "] :"
                    cell_text3 = cell3.text_frame.add_paragraph()  
                    self.__edit_paragraph(cell_text3, text_content=responsibilities_ref, font_size=11, font_name='Calibri', color=random_direct_cause_event_color, bold=True, alignment=PP_ALIGN.CENTER)

                    responsibles = improvement_action['responsible']
                    responsibles = responsibles.replace("\r", "").replace("\n","")

                    cell_text3 = cell3.text_frame.add_paragraph()  
                    self.__edit_paragraph(cell_text3, text_content=responsibles, font_size=11, font_name='Calibri', alignment=PP_ALIGN.CENTER)

                    cell_text3 = cell3.text_frame.add_paragraph()
                    cell_text3 = cell3.text_frame.add_paragraph()


                    cell4 = table.cell(1, 4)
                    cell_text4 = cell4.text_frame.add_paragraph()  
                    self.__edit_paragraph(cell_text4, text_content=responsibilities_ref, font_size=11, font_name='Calibri', color=random_direct_cause_event_color, bold=True, alignment=PP_ALIGN.CENTER)

                    date_obj = datetime.fromisoformat(improvement_action['deadLine'])
                    incident_report_edition_date = date_obj.strftime('%d/%m/%Y')

                    cell_text4 = cell4.text_frame.add_paragraph()  
                    self.__edit_paragraph(cell_text4, text_content=incident_report_edition_date, font_size=11, font_name='Calibri', alignment=PP_ALIGN.CENTER)

                    cell_text4 = cell4.text_frame.add_paragraph()
                    cell_text4 = cell4.text_frame.add_paragraph()


    def set_resume_slide(self, slide_title):

        # Ajouter une diapositive vide (la septième diapositive)
        slide_layout = self.prs.slide_layouts[6]  # 6 est une diapositive vide
        slide9 = self.prs.slides.add_slide(slide_layout)

        title_box = self.add_textbox_to_slide(slide9.slide_id, left=0, top=0, width=10, height=1)
        self.edit_textbox(title_box, text=slide_title, font_size=30, font_name='Calibri', color=bright_red, alignment=PP_ALIGN.CENTER, bold=True)

        self.add_textbox_to_slide(slide9.slide_id, left=0.25, top=1, width=9.5, height=6, word_wrap=True)


    def set_appendix(self, text):
        title_slide_layout = self.prs.slide_layouts[0]
        slide5 = self.prs.slides.add_slide(title_slide_layout)

        title_shape = slide5.shapes.placeholders[0]

        self.__edit_paragraph(title_shape.text_frame.paragraphs[0], text_content=text, color=bright_red, bold=True)


    def set_end_slide(self, end_title):
        blank_slide_layout = self.prs.slide_layouts[6]  # 6 est une diapositive vide
        slide6 = self.prs.slides.add_slide(blank_slide_layout)

        title_slide6_shape = self.add_textbox_to_slide(slide6.slide_id, left=2.8, top=2.75, width=4.5, height=2, word_wrap=False)
        self.edit_textbox(title_slide6_shape, text=end_title, font_size=115, color=sky_blue, bold=True)



